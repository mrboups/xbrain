"""MCP Deck sidecar — PowerPoint generator via python-pptx.

Standalone FastMCP server (streamable-http transport, port 8103).
Generates .pptx files from structured section data, uploads to MinIO,
indexes the asset in memory-api with full tagging contract.

Tools:
  deck_create(title, sections) -> {url: str, deck_id: str, memory_item_id: str}
  deck_update(deck_id, sections) -> {url: str, deck_id: str, version: int}

IMPORTANT: Do NOT mount this inside a parent FastAPI app (issue #1367).
Single worker mandatory (FastMCP session state is in-memory per process).
"""
from __future__ import annotations

import io
import os
import uuid
from typing import Any

import httpx
import structlog
from mcp.server.fastmcp import FastMCP

log = structlog.get_logger(__name__)

# ── Config ───────────────────────────────────────────────────────────────────
MINIO_ENDPOINT = os.environ.get("MINIO_ENDPOINT", "minio:9000")
MINIO_ACCESS_KEY = os.environ.get("MINIO_ACCESS_KEY", "")
MINIO_SECRET_KEY = os.environ.get("MINIO_SECRET_KEY", "")
MINIO_BUCKET = os.environ.get("MINIO_BUCKET_DECKS", "xbrain-decks")
MINIO_USE_SSL = os.environ.get("MINIO_USE_SSL", "false").lower() == "true"
MEMORY_API_URL = os.environ.get("MEMORY_API_URL", "http://memory-api:8000")
BRIDGE_SHARED_SECRET = os.environ.get("BRIDGE_SHARED_SECRET", "")
HOST = os.environ.get("FASTMCP_HOST", "0.0.0.0")
PORT = int(os.environ.get("FASTMCP_PORT", "8103"))

# DoS guard limits (T-04-07-SEC-02)
MAX_SECTIONS = 50
MAX_BULLETS_PER_SECTION = 20
MAX_BULLET_CHARS = 500
MAX_TITLE_CHARS = 200

mcp = FastMCP("xbrain-deck", host=HOST, port=PORT)


# ── MinIO helpers ─────────────────────────────────────────────────────────────

def _get_s3_client():
    import boto3
    protocol = "https" if MINIO_USE_SSL else "http"
    return boto3.client(
        "s3",
        endpoint_url=f"{protocol}://{MINIO_ENDPOINT}",
        aws_access_key_id=MINIO_ACCESS_KEY,
        aws_secret_access_key=MINIO_SECRET_KEY,
        region_name="us-east-1",  # MinIO ignores region but boto3 requires it
    )


def _ensure_bucket(s3) -> None:
    """Create bucket if it does not exist."""
    try:
        s3.head_bucket(Bucket=MINIO_BUCKET)
    except Exception:
        s3.create_bucket(Bucket=MINIO_BUCKET)
        log.info("deck.bucket_created", bucket=MINIO_BUCKET)


def _upload_pptx(deck_id: str, pptx_bytes: bytes, version: int = 1) -> str:
    """Upload .pptx to MinIO and return presigned URL (1h expiry)."""
    s3 = _get_s3_client()
    _ensure_bucket(s3)
    key = f"decks/{deck_id}/v{version}.pptx"
    s3.put_object(
        Bucket=MINIO_BUCKET,
        Key=key,
        Body=pptx_bytes,
        ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    url = s3.generate_presigned_url(
        "get_object",
        Params={"Bucket": MINIO_BUCKET, "Key": key},
        ExpiresIn=3600,
    )
    log.info("deck.uploaded", deck_id=deck_id, version=version, key=key)
    return url


def _get_current_version(deck_id: str) -> int:
    """Return the current highest version number for a deck. Returns 0 if none."""
    s3 = _get_s3_client()
    try:
        resp = s3.list_objects_v2(Bucket=MINIO_BUCKET, Prefix=f"decks/{deck_id}/v")
        objects = resp.get("Contents", [])
        versions = []
        for obj in objects:
            key = obj["Key"]
            # Extract version number from "decks/{id}/v{n}.pptx"
            parts = key.split("/v")
            if len(parts) == 2:
                try:
                    versions.append(int(parts[1].replace(".pptx", "")))
                except ValueError:
                    pass
        return max(versions) if versions else 0
    except Exception:
        return 0


# ── python-pptx helpers ───────────────────────────────────────────────────────

def _generate_pptx(title: str, sections: list[dict]) -> bytes:
    """Generate a .pptx file from title and sections.

    sections: list of {heading: str, bullets: list[str]}
    Returns raw bytes of the .pptx file.

    Uses python-pptx default blank presentation.
    Each section becomes one slide: title=heading, body=bullets.
    First slide is a title slide with the deck title.
    """
    from pptx import Presentation

    prs = Presentation()

    # Title slide (layout 0 = title slide)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title[:MAX_TITLE_CHARS]
    # Subtitle placeholder (index 1) — optional
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"{len(sections)} section(s)"

    # Content slides (layout 1 = title + content)
    content_layout = prs.slide_layouts[1]
    for section in sections:
        heading = str(section.get("heading", ""))[:MAX_TITLE_CHARS]
        bullets = section.get("bullets", [])[:MAX_BULLETS_PER_SECTION]
        if not heading and not bullets:
            continue
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = heading
        # Add bullets to the body placeholder
        tf = s.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = str(bullet)[:MAX_BULLET_CHARS]
            p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── memory-api indexing ───────────────────────────────────────────────────────

def _mint_bridge_jwt(team_scope: str) -> str:
    import time
    from joserfc import jwt as jose_jwt
    from joserfc.jwk import OctKey
    now = int(time.time())
    claims = {
        "iss": "mcp-deck",
        "sub": "mcp-deck",
        "scope": "bridge",
        "team_scope": team_scope,
        "iat": now,
        "exp": now + 300,
    }
    key = OctKey.import_key(
        BRIDGE_SHARED_SECRET.encode() if isinstance(BRIDGE_SHARED_SECRET, str) else BRIDGE_SHARED_SECRET
    )
    token = jose_jwt.encode({"alg": "HS256"}, claims, key)
    return token


def _index_in_memory_api(
    deck_id: str, title: str, url: str, team_scope: str, user_sub: str
) -> str | None:
    """Index the deck asset in memory-api. Returns memory_item_id or None on error."""
    if not BRIDGE_SHARED_SECRET:
        log.warning("deck.memory_index_skipped", reason="BRIDGE_SHARED_SECRET not set")
        return None
    import datetime
    jwt = _mint_bridge_jwt(team_scope)
    now_iso = datetime.datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
    payload = {
        "item": {
            "id": deck_id,
            "content": f"Pitch deck: {title}",
            "team_scope": team_scope,
            "project_scope": None,
            "visibility": "team",
            "confidence": 1.0,
            "truth_level": "WORKING",
            "source": f"mcp:deck:{deck_id}",
            "validation_status": "pending",
            "created_at": now_iso,
            "updated_at": now_iso,
            "metadata": {
                "deck_id": deck_id,
                "minio_url": url,
                "title": title,
                "generated_by": user_sub,
            },
        }
    }
    try:
        r = httpx.post(
            f"{MEMORY_API_URL}/v1/memory/upsert",
            json=payload,
            headers={
                "Authorization": f"Bearer {jwt}",
                "X-Team-Scope": team_scope,
                "Content-Type": "application/json",
            },
            timeout=15.0,
        )
        r.raise_for_status()
        return r.json().get("id")
    except Exception as exc:
        log.warning("deck.memory_index_failed", deck_id=deck_id, error=str(exc))
        return None


# ── Input validation (T-04-07-SEC-02 DoS guard) ──────────────────────────────

def _validate_sections(sections: list) -> None:
    """Enforce size limits to prevent DoS via enormous payloads."""
    if len(sections) > MAX_SECTIONS:
        raise ValueError(f"max {MAX_SECTIONS} sections allowed, got {len(sections)}")
    for i, section in enumerate(sections):
        if not isinstance(section, dict):
            raise ValueError(f"section[{i}] must be a dict with 'heading' and 'bullets'")
        bullets = section.get("bullets", [])
        if len(bullets) > MAX_BULLETS_PER_SECTION:
            raise ValueError(
                f"section[{i}] has {len(bullets)} bullets; max is {MAX_BULLETS_PER_SECTION}"
            )


# ── MCP Tools ─────────────────────────────────────────────────────────────────

@mcp.tool()
async def deck_create(
    title: str,
    sections: list,
    team_scope: str = "default",
    user_sub: str = "unknown",
) -> dict:
    """Create a PowerPoint presentation from structured section data.

    Args:
        title: The deck title (shown on title slide). Max 200 chars.
        sections: List of sections, each with "heading" (str) and "bullets" (list[str]).
                  Max 50 sections, 20 bullets per section, 500 chars per bullet.
                  Example: [{"heading": "Problem", "bullets": ["Pain 1", "Pain 2"]}]
        team_scope: Team scope for storage and indexing (injected by gateway).
        user_sub: User sub for audit trail (injected by gateway).

    Returns:
        dict with "url" (presigned MinIO URL, 1h expiry), "deck_id", "memory_item_id",
        and "version" (always 1 for new decks).
    """
    # DoS guard
    _validate_sections(sections)

    deck_id = str(uuid.uuid4())
    log.info(
        "deck.create_start",
        deck_id=deck_id,
        title=title[:50],
        sections=len(sections),
        team=team_scope,
    )

    pptx_bytes = _generate_pptx(title, sections)
    url = _upload_pptx(deck_id, pptx_bytes, version=1)
    memory_item_id = _index_in_memory_api(deck_id, title, url, team_scope, user_sub)

    log.info("deck.created", deck_id=deck_id, memory_item_id=memory_item_id)
    return {"url": url, "deck_id": deck_id, "memory_item_id": memory_item_id, "version": 1}


@mcp.tool()
async def deck_update(
    deck_id: str,
    sections: list,
    title: str = "",
    team_scope: str = "default",
    user_sub: str = "unknown",
) -> dict:
    """Update an existing deck by regenerating it with new sections.

    The previous version is preserved as decks/{deck_id}/v{n}.pptx in MinIO.
    Returns a new presigned URL for the updated version.

    Args:
        deck_id: The deck ID returned by deck_create.
        sections: New sections to replace the existing content.
        title: Optional updated title (defaults to "Deck {deck_id}" if omitted).
        team_scope: Team scope (injected by gateway).
        user_sub: User sub (injected by gateway).

    Returns:
        dict with "url" (new presigned URL), "deck_id", "version" (new version number).
    """
    # DoS guard
    _validate_sections(sections)

    log.info(
        "deck.update_start",
        deck_id=deck_id,
        sections=len(sections),
        team=team_scope,
    )

    current_version = _get_current_version(deck_id)
    new_version = current_version + 1

    deck_title = title if title else f"Deck {deck_id}"
    pptx_bytes = _generate_pptx(deck_title, sections)
    url = _upload_pptx(deck_id, pptx_bytes, version=new_version)

    log.info("deck.updated", deck_id=deck_id, version=new_version)
    return {"url": url, "deck_id": deck_id, "version": new_version}


if __name__ == "__main__":
    # Single worker — critical: FastMCP session state is in-memory per process.
    # Multi-worker mode causes session 404s (issue #658).
    # Transport streamable-http binds to /mcp endpoint on the specified port.
    mcp.run(transport="streamable-http")
